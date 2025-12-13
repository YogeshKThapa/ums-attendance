from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define some basic slide layouts
    # 0 = Title Slide
    # 1 = Title and Content
    # 6 = Blank

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text

    def add_content_slide(title_text, content_items):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Access the body placeholder
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        # Add content
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.size = Pt(20)

    # --- SLIDE 1 ---
    add_title_slide("UMS Attendance Pro", "Revolutionizing the Student Attendance Experience\nOverview & Features")

    # --- SLIDE 2 ---
    add_content_slide("The Problem", [
        "❌ Slow Access: Official portal is sluggish.",
        "❌ Poor UX: Raw tables are hard to read.",
        "❌ No Insights: Manual calculation for 'Safe to Bunk'.",
        "❌ Fragmented Data: No 'Overall Semester' view."
    ])

    # --- SLIDE 3 ---
    add_content_slide("The Solution: UMS Attendance Pro", [
        "A High-Performance Web Application (PWA) wrapper.",
        "",
        "⚡ Fast: Caches data for offline access.",
        "🧠 Smart: Auto-calculates metrics.",
        "📊 Unified: Aggregates all months.",
        "🏆 Social: Global Leaderboard."
    ])

    # --- SLIDE 4 ---
    add_content_slide("Key Features (1/2)", [
        "1. Smart 'Safe-to-Bunk' Meter",
        "   - Visual progress bars (Green/Red).",
        "   - Instant actionable feedback.",
        "",
        "2. Overall Semester View",
        "   - Fetches all months in parallel.",
        "   - Aggregates data for accurate exam eligibility."
    ])

    # --- SLIDE 5 ---
    add_content_slide("Key Features (2/2)", [
        "3. Session Persistence",
        "   - Auto-recovers expired sessions.",
        "   - Reduces CAPTCHA fatigue.",
        "",
        "4. Automatic Leaderboard (Social)",
        "   - Gamified ranking system.",
        "   - Background automatic updates.",
        "   - Privacy-focused implementation."
    ])

    # --- SLIDE 6 ---
    add_content_slide("Technology Stack", [
        "Frontend:",
        "• React (Vite) + PWA",
        "",
        "Backend:",
        "• Python (Flask)",
        "• MongoDB Atlas (Leaderboard)",
        "• BeautifulSoup4 (Scraping)",
        "",
        "Deployment:",
        "• Vercel (Frontend)",
        "• Render (Backend)"
    ])

    # --- SLIDE 7 ---
    add_content_slide("Future Roadmap", [
        "🤖 AI Bunk Planner Simulator",
        "🔔 Push Notifications",
        "📱 Native Mobile App (React Native)",
        "📅 Google Calendar Integration"
    ])

    # Save
    output_path = r"C:\Users\kumar\.gemini\antigravity\brain\f694cb56-12d1-4c27-94e7-4a49a36cab5e\UMS_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
